from pptx import Presentation
from pptx.util import Inches, Pt
import os, time
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from TM1py import TM1Service
import configparser

# ========================================================================
# Region - Definition
## Read information from config.ini (in currenct path)
config = configparser.ConfigParser()
config.read('config.ini', encoding='utf-8')

## PPT_Maker
SCREENSHOT_DIR  = config['PPT_Maker']['ppt_screenshot_path']
OUTPUT_PPTX     = config['PPT_Maker']['ppt_path']     # PPT path
MASTER_PPT      = config['PPT_Maker']['master_ppt_path'] 
PPT_UX_APP_MDX  = config['PPT_Maker']['ppt_ux_app_mdx']        # 想要製作投影片的的UX Apps, 將透過此UX Apps結果去尋找圖檔 (define by MDX)

COVER_PAGE_IDX  = int( config['PPT_Maker']['cover_page_idx'])
FINAL_PAGE_IDX  = int( config['PPT_Maker']['final_page_idx'])
BLANK_PAGE_IDX  = int( config['PPT_Maker']['blank_page_idx'])

TITLE_LEFT      = float( config['PPT_Maker']['title_left'])
TITLE_TOP       = float( config['PPT_Maker']['title_top'])
TITLE_WIDTH     = float( config['PPT_Maker']['title_width'])
TITLE_HEIGHT    = float( config['PPT_Maker']['title_height'])
CONTENT_LEFT    = float( config['PPT_Maker']['content_left'])
CONTENT_TOP     = float( config['PPT_Maker']['content_top'])
CONTENT_WIDTH   = float( config['PPT_Maker']['content_width'])
CONTENT_HEIGHT  = float( config['PPT_Maker']['content_height'])
PICTURE_LEFT    = float( config['PPT_Maker']['picture_left'])
PICTURE_TOP     = float( config['PPT_Maker']['picture_top'])
PICTURE_WIDTH   = float( config['PPT_Maker']['picture_width'])
PICTURE_HEIGHT  = float( config['PPT_Maker']['picture_height'])
BOX_COLOR       = config['PPT_Maker']['box_color']
B_R, B_G, B_B   = map(int, BOX_COLOR.split(","))
TEXT_COLOR      = config['PPT_Maker']['text_color']
T_R, T_G, T_B   = map(int, TEXT_COLOR.split(","))
TEXT_SIZE       = float( config['PPT_Maker']['text_size'])

## UX_CS
ADDRESS         = config['UX_CS']['address']         # Content Store Address
PORT            = config['UX_CS']['port']            # Content Store Port
NAMESPACE       = config['UX_CS']['namespace']       # Content Store Namespace
SYS_USERNAME    = config['UX_CS']['system_username'] # Content Store Login System Account
SYS_PASSWORD    = config['UX_CS']['system_password'] # Content Store Login System Account Password
SSL             = config['UX_CS']['ssl']             # Content Store Using SSL?

## UX_CS Objects
dimension_name      = "}APQ UX App"
attr_name_title     = "Page Title"
attr_name_func      = "Page Description"
attr_name_idx       = "Page Index"
attr_name_default   = "Code and Description"

# EndRegion - Definition
# ========================================================================

# ========================================================================
# Region - Get UX App List
# === TM1 Connection - TM1 連線 ===
tm1 = TM1Service(address=ADDRESS, port=PORT, user=SYS_USERNAME, password=SYS_PASSWORD, ssl=SSL, namespace=NAMESPACE)
mdx = 'Order( ' + PPT_UX_APP_MDX + ' , [' + dimension_name + '].[' + dimension_name + '].CurrentMember.Properties("' + attr_name_idx + '"), ASC)'
ux_apps = tm1.elements.execute_set_mdx_element_names( mdx)
 
# EndRegion - Get UX App List
# ========================================================================

# ========================================================================
# Region - PPT Definition
# PPT Master Slides Template
prs = Presentation(MASTER_PPT)

# Page definition
cover_slide_layout = prs.slide_layouts[COVER_PAGE_IDX]
final_slide_layout = prs.slide_layouts[FINAL_PAGE_IDX]
blank_slide_layout = prs.slide_layouts[BLANK_PAGE_IDX]

# Create Cover Page - 建立封面投影片
cover_slide = prs.slides.add_slide(cover_slide_layout)

# EndRegion - PPT Definition
# ========================================================================

# ========================================================================
# Region - Create Page Title/ Page Description dictionary
title_dict = tm1.elements.get_attribute_of_elements(
                dimension_name='}APQ UX App', 
                hierarchy_name='}APQ UX App',
                attribute=attr_name_title,
            )
            
desc_dict = tm1.elements.get_attribute_of_elements(
                dimension_name='}APQ UX App', 
                hierarchy_name='}APQ UX App',
                attribute=attr_name_func,
            )

code_desc_dict = tm1.elements.get_attribute_of_elements(
                dimension_name='}APQ UX App', 
                hierarchy_name='}APQ UX App',
                attribute=attr_name_default,
            )
# EndRegion - Create Page Title/ Page Description dictionary
# ========================================================================

# ========================================================================
# Region - Insert Slides by screenshot
# === Loop UX App element by sorting ===
for app_name in ux_apps:
    filename = app_name + '.png'
    if app_name:
        element = app_name
        # Get the Page Title attribute from TM1/ 從 TM1 抓取 Page Title attribute
        ppt_title = title_dict.get(element, element)
        if ppt_title == element:
            ppt_title = code_desc_dict.get(element, element)  # If not caught, use the Code and Description attribute as title/ 如果沒抓到，用 Code and Description attribute當標題
            
        # Get the Page Description attribute from TM1/ 從 TM1 抓取 Page Description attribute
        func_desc = desc_dict.get(element, element)
        if not func_desc:
            func_desc = element  # If not caught, use the element name as content/ 如果沒抓到，用 element 名稱當內容
        
        # Create Blank Page - 建立空白投影片
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add a title/ 加入標題
        t_left      = Inches(TITLE_LEFT)
        t_top       = Inches(TITLE_TOP)
        t_width     = Inches(TITLE_WIDTH)
        t_height    = Inches(TITLE_HEIGHT)
        t_box = slide.shapes.add_textbox(
                        t_left      # Title X-axis/ Title X軸
                        , t_top     # Title Y-axis/ Title Y軸
                        , t_width   # Title Width/  Title 寬度
                        , t_height  # Title Height/ Title 高度
                        )
        tf = t_box.text_frame
        p = tf.add_paragraph()
        p.text = ppt_title
        p.font.size = Pt(28)
        
        
        # 插入綠色區塊（矩形）
        c_left      = Inches(CONTENT_LEFT)
        c_top       = Inches(CONTENT_TOP)
        c_width     = Inches(CONTENT_WIDTH)
        c_height    = Inches(CONTENT_HEIGHT)
        c_box = slide.shapes.add_shape(
                        MSO_AUTO_SHAPE_TYPE.RECTANGLE
                        , c_left      # Content X-axis/ Content X軸
                        , c_top       # Content Y-axis/ Content Y軸
                        , c_width     # Content Width/  Content 寬度
                        , c_height    # Content Height/ Content 高度
                        )
        fill = c_box.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(B_R, B_G, B_B)  # Box Color
        c_box.line.fill.background()  # Remove border line/ 移除邊框線
        
        tf = c_box.text_frame
        tf.clear()  # Clear the default text/ 清掉預設文字
        
        # Set text to upper left/ 設定文字靠左上
        tf.margin_left = Pt(5)
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"Function: {func_desc}"
        run.font.size = Pt(TEXT_SIZE)  
        run.font.bold = False
        run.font.color.rgb = RGBColor(T_R, T_G, T_B)

        # Add pictures/ 加入圖片
        img_path    = os.path.join(SCREENSHOT_DIR, filename)
        p_left      = Inches(PICTURE_LEFT)
        p_top       = Inches(PICTURE_TOP)
        p_width     = Inches(PICTURE_WIDTH)
        p_height    = Inches(PICTURE_HEIGHT)
        slide.shapes.add_picture(
            img_path
            , p_left    # Picture X-axis/ 圖片X軸
            , p_top     # Picture Y-axis/ 圖片Y軸
            , p_width   # Picture Width/  圖片寬度
            , p_height  # Picture Height/ 圖片高度
            )

# EndRegion - Insert Slides by screenshot
# ========================================================================

# ========================================================================
# Region - Final Page
# Create Final Page/ 建立尾頁
final_slide = prs.slides.add_slide(final_slide_layout)

# EndRegion - Final Page
# ========================================================================

# === Save PPT/ 儲存 PPT ===
prs.save(OUTPUT_PPTX)
print(f"✅ Successfully generated PowerPoint：{OUTPUT_PPTX}")

# === Sleep 1.5 sec for check result/ 等待1.5秒, 查看結果 ===
time.sleep(1.5)
